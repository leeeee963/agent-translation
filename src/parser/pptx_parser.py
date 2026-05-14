"""PPTX parser: per-paragraph block model.

Every PPTX paragraph that carries translatable text becomes its own
ContentBlock. The LLM returns one [[BLOCK:id]] marker per paragraph, so
- a missing block triggers the existing retry / single-block-fallback path
  in TranslatorAgent (no silent line drops)
- rebuild looks up by (text_frame_group, para_index) and writes directly
  into the matching paragraph (no \\n splitting heuristics)

Block ids
---------
- Text frame paragraph:   slide{N}_shape{M}[_shape{X}...]_p{K}
- Table cell paragraph:   slide{N}_shape{M}_r{R}c{C}_p{K}
- Notes paragraph:        slide{N}_note_p{K}
- Chart title:            slide{N}_shape{M}_chart_title
- Chart categories:       slide{N}_shape{M}_chart_cats
- Chart series:           slide{N}_shape{M}_chart_s{I}
- SmartArt diagram item:  slide{N}_dgm{I}    (handled by pptx_diagram)

Each per-paragraph block carries `text_frame_group` (the parent's stable id)
and `para_index` (its position in tf.paragraphs) in metadata.
"""

from __future__ import annotations

import logging
import os
import shutil
import uuid
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

from src.models.content import BlockType, ContentBlock, FileMeta, ParsedFile
from src.parser.base import BaseParser
from src.parser.pptx_diagram import build_diagram_map, parse_diagram, rebuild_diagrams
from src.parser.pptx_text import (
    _flatten_for_paragraph,
    get_para_dominant_fmt,
    warn_overflow,
    write_para_with_fmt,
    write_paragraphs_by_index,
)
from src.utils.layout_fixer import adjust_runs_font_size, enable_autofit

logger = logging.getLogger(__name__)

_HEADING_FONT_SIZE_THRESHOLD = Pt(24)


class PptxParser(BaseParser):
    """Parser for Microsoft PowerPoint (.pptx) files."""

    EXTENSIONS = {".pptx"}

    def can_handle(self, file_path: str) -> bool:
        return Path(file_path).suffix.lower() in self.EXTENSIONS

    # ------------------------------------------------------------------
    # parse
    # ------------------------------------------------------------------
    def parse(self, file_path: str) -> ParsedFile:
        prs = Presentation(file_path)
        blocks: list[ContentBlock] = []

        diagram_map = build_diagram_map(file_path)

        for slide_idx, slide in enumerate(prs.slides):
            slide_blocks = self._parse_shape_tree(slide.shapes, slide_idx)
            blocks.extend(slide_blocks)

            slide_num = slide_idx + 1
            if slide_num in diagram_map:
                dgm_blocks = parse_diagram(
                    file_path, diagram_map[slide_num], slide_idx
                )
                blocks.extend(dgm_blocks)

            if slide.has_notes_slide:
                blocks.extend(self._parse_notes(slide, slide_idx))

        total_words = sum(len(b.source_text.split()) for b in blocks)

        return ParsedFile(
            meta=FileMeta(
                original_name=os.path.basename(file_path),
                file_type="pptx",
                word_count=total_words,
            ),
            blocks=blocks,
            format_template=str(Path(file_path).resolve()),
        )

    # ── shape tree (text frames, tables, groups) ─────────────────────

    def _parse_shape_tree(
        self, shapes, slide_idx: int, prefix: str = ""
    ) -> list[ContentBlock]:
        blocks: list[ContentBlock] = []
        for shape_idx, shape in enumerate(shapes):
            shape_id = (
                f"{prefix}shape{shape_idx}" if prefix else f"shape{shape_idx}"
            )
            block_id_base = f"slide{slide_idx}_{shape_id}"
            blocks_before = len(blocks)

            # GroupShape — recurse
            if shape.shape_type is not None and shape.shape_type == 6:
                try:
                    blocks.extend(
                        self._parse_shape_tree(
                            shape.shapes, slide_idx, prefix=f"{shape_id}_"
                        )
                    )
                except Exception:
                    logger.debug("Could not iterate group shape %s", block_id_base)
                continue

            # Chart — extract title, categories, series names (single-string blocks)
            if hasattr(shape, "has_chart") and shape.has_chart:
                chart_blocks = self._parse_chart(shape.chart, slide_idx, block_id_base)
                blocks.extend(chart_blocks)

            # Table — per-paragraph blocks per cell
            if shape.has_table:
                blocks.extend(
                    self._parse_table(shape.table, slide_idx, block_id_base)
                )

            # Text frame — per-paragraph blocks
            if shape.has_text_frame:
                blocks.extend(
                    self._parse_text_frame(
                        shape.text_frame,
                        slide_idx,
                        shape_idx,
                        block_id_base,
                    )
                )
            elif not shape.has_table and not (
                hasattr(shape, "has_chart") and shape.has_chart
            ):
                # Fallback: shapes without text_frame but with .text — single block,
                # no per-para breakdown possible (no paragraph structure exposed).
                try:
                    text = shape.text.strip() if hasattr(shape, "text") else ""
                    if text:
                        blocks.append(
                            ContentBlock(
                                id=block_id_base,
                                type=BlockType.PARAGRAPH,
                                source_text=text,
                                metadata={
                                    "slide_index": slide_idx,
                                    "shape_index": shape_idx,
                                    "shape_kind": "fallback_text",
                                },
                            )
                        )
                except Exception:
                    pass

            if len(blocks) == blocks_before:
                logger.debug(
                    "Shape %s (type=%s) produced no content blocks",
                    block_id_base,
                    getattr(shape, "shape_type", "unknown"),
                )

        return blocks

    def _parse_text_frame(
        self,
        tf,
        slide_idx: int,
        shape_idx: int,
        block_id_base: str,
    ) -> list[ContentBlock]:
        blocks: list[ContentBlock] = []
        paras_info = [get_para_dominant_fmt(para) for para in tf.paragraphs]
        total_paras = len(paras_info)

        for para_idx, pi in enumerate(paras_info):
            raw = pi.get("translatable_text") or ""
            text = raw.strip()
            if not text:
                continue  # spacer or icon-only — no translatable content

            # Preserve any whitespace surrounding the translatable text so we
            # can re-apply it at rebuild (matters for icon-prefixed paragraphs
            # like "<icon>favorite</icon> liked item" where the leading space
            # is the visual separator).
            lead_ws = raw[: len(raw) - len(raw.lstrip())]
            trail_ws = raw[len(raw.rstrip()) :]

            font_size = pi.get("font_size") or 0
            block_type = (
                BlockType.HEADING
                if font_size >= _HEADING_FONT_SIZE_THRESHOLD
                else BlockType.PARAGRAPH
            )

            blocks.append(
                ContentBlock(
                    id=f"{block_id_base}_p{para_idx}",
                    type=block_type,
                    source_text=text,
                    metadata={
                        "slide_index": slide_idx,
                        "shape_index": shape_idx,
                        "shape_kind": "text_frame",
                        "text_frame_group": block_id_base,
                        "group_kind": "text_frame",
                        "para_index": para_idx,
                        "total_paras": total_paras,
                        "para_format": pi,
                        "lead_ws": lead_ws,
                        "trail_ws": trail_ws,
                    },
                )
            )
        return blocks

    def _parse_table(
        self, table, slide_idx: int, block_id_base: str
    ) -> list[ContentBlock]:
        blocks: list[ContentBlock] = []
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                cell_group = f"{block_id_base}_r{row_idx}c{col_idx}"
                paras_info = [
                    get_para_dominant_fmt(p)
                    for p in cell.text_frame.paragraphs
                ]
                total_paras = len(paras_info)
                for para_idx, pi in enumerate(paras_info):
                    raw = pi.get("translatable_text") or ""
                    text = raw.strip()
                    if not text:
                        continue
                    lead_ws = raw[: len(raw) - len(raw.lstrip())]
                    trail_ws = raw[len(raw.rstrip()) :]
                    blocks.append(
                        ContentBlock(
                            id=f"{cell_group}_p{para_idx}",
                            type=BlockType.PARAGRAPH,
                            source_text=text,
                            metadata={
                                "slide_index": slide_idx,
                                "shape_kind": "table_cell",
                                "text_frame_group": cell_group,
                                "group_kind": "table_cell",
                                "para_index": para_idx,
                                "total_paras": total_paras,
                                "para_format": pi,
                                "table_block": block_id_base,
                                "row": row_idx,
                                "col": col_idx,
                                "lead_ws": lead_ws,
                                "trail_ws": trail_ws,
                            },
                        )
                    )
        return blocks

    def _parse_notes(self, slide, slide_idx: int) -> list[ContentBlock]:
        blocks: list[ContentBlock] = []
        notes_tf = slide.notes_slide.notes_text_frame
        notes_group = f"slide{slide_idx}_note"
        paras_info = [get_para_dominant_fmt(p) for p in notes_tf.paragraphs]
        total_paras = len(paras_info)
        for para_idx, pi in enumerate(paras_info):
            raw = pi.get("translatable_text") or ""
            text = raw.strip()
            if not text:
                continue
            lead_ws = raw[: len(raw) - len(raw.lstrip())]
            trail_ws = raw[len(raw.rstrip()) :]
            blocks.append(
                ContentBlock(
                    id=f"{notes_group}_p{para_idx}",
                    type=BlockType.SLIDE_NOTE,
                    source_text=text,
                    metadata={
                        "slide_index": slide_idx,
                        "text_frame_group": notes_group,
                        "group_kind": "notes",
                        "para_index": para_idx,
                        "total_paras": total_paras,
                        "para_format": pi,
                        "lead_ws": lead_ws,
                        "trail_ws": trail_ws,
                    },
                )
            )
        return blocks

    def _parse_chart(
        self, chart, slide_idx: int, block_id_base: str
    ) -> list[ContentBlock]:
        """Extract translatable text from chart: title, category labels, series names."""
        blocks: list[ContentBlock] = []

        # Chart title
        try:
            if chart.has_title and chart.chart_title.has_text_frame:
                title_text = chart.chart_title.text_frame.text.strip()
                if title_text:
                    blocks.append(
                        ContentBlock(
                            id=f"{block_id_base}_chart_title",
                            type=BlockType.HEADING,
                            source_text=title_text,
                            metadata={
                                "slide_index": slide_idx,
                                "shape_kind": "chart_title",
                            },
                        )
                    )
        except Exception:
            logger.debug("Could not extract chart title from %s", block_id_base)

        # Category axis labels
        try:
            plot = chart.plots[0]
            categories = [str(cat) for cat in plot.categories if str(cat).strip()]
            if categories:
                cat_text = " | ".join(categories)
                blocks.append(
                    ContentBlock(
                        id=f"{block_id_base}_chart_cats",
                        type=BlockType.PARAGRAPH,
                        source_text=cat_text,
                        metadata={
                            "slide_index": slide_idx,
                            "shape_kind": "chart_categories",
                            "category_count": len(categories),
                        },
                    )
                )
        except Exception:
            logger.debug("Could not extract chart categories from %s", block_id_base)

        # Series names
        try:
            for s_idx, series in enumerate(chart.series):
                try:
                    name = series.name if hasattr(series, "name") else None
                    if name and str(name).strip():
                        blocks.append(
                            ContentBlock(
                                id=f"{block_id_base}_chart_s{s_idx}",
                                type=BlockType.PARAGRAPH,
                                source_text=str(name).strip(),
                                metadata={
                                    "slide_index": slide_idx,
                                    "shape_kind": "chart_series",
                                    "series_index": s_idx,
                                },
                            )
                        )
                except Exception:
                    pass
        except Exception:
            logger.debug("Could not extract chart series from %s", block_id_base)

        return blocks

    # ------------------------------------------------------------------
    # rebuild
    # ------------------------------------------------------------------
    def rebuild(self, parsed_file: ParsedFile, output_path: str) -> str:
        original_path = parsed_file.format_template
        if not original_path or not os.path.exists(original_path):
            raise FileNotFoundError(
                f"Original PPTX not found: {original_path}"
            )

        tmp_path = f"{output_path}.tmp_{uuid.uuid4().hex[:8]}.pptx"
        shutil.copy2(original_path, tmp_path)
        prs = Presentation(tmp_path)

        block_map = self._build_block_map(parsed_file)
        group_map = self._build_group_map(block_map)

        for slide_idx, slide in enumerate(prs.slides):
            self._rebuild_shape_tree(slide.shapes, slide_idx, block_map, group_map)
            self._rebuild_notes(slide, slide_idx, group_map)

        prs.save(tmp_path)

        diagram_blocks = [
            b for b in parsed_file.blocks
            if b.metadata.get("shape_kind") == "diagram"
        ]
        if diagram_blocks:
            rebuild_diagrams(tmp_path, output_path, block_map)
        else:
            shutil.move(tmp_path, output_path)

        if os.path.exists(tmp_path):
            os.remove(tmp_path)

        logger.info("Rebuilt PPTX saved to %s", output_path)
        return output_path

    @staticmethod
    def _build_block_map(parsed_file: ParsedFile) -> dict[str, ContentBlock]:
        return {b.id: b for b in parsed_file.blocks}

    @staticmethod
    def _build_group_map(
        block_map: dict[str, ContentBlock],
    ) -> dict[str, dict[int, ContentBlock]]:
        """Group per-paragraph blocks by their parent text_frame_group.

        Returns ``{group_id: {para_index: block}}`` for fast rebuild lookup.
        Blocks without ``text_frame_group`` (charts, fallback shapes,
        diagrams) are ignored here — they take a different rebuild path.
        """
        group_map: dict[str, dict[int, ContentBlock]] = {}
        for block in block_map.values():
            group = block.metadata.get("text_frame_group")
            if not group:
                continue
            para_idx = block.metadata.get("para_index")
            if para_idx is None:
                continue
            group_map.setdefault(group, {})[int(para_idx)] = block
        return group_map

    def _rebuild_shape_tree(
        self,
        shapes,
        slide_idx: int,
        block_map: dict[str, ContentBlock],
        group_map: dict[str, dict[int, ContentBlock]],
        prefix: str = "",
    ) -> None:
        for shape_idx, shape in enumerate(shapes):
            shape_id = (
                f"{prefix}shape{shape_idx}" if prefix else f"shape{shape_idx}"
            )
            block_id_base = f"slide{slide_idx}_{shape_id}"

            if shape.shape_type is not None and shape.shape_type == 6:
                try:
                    self._rebuild_shape_tree(
                        shape.shapes, slide_idx, block_map, group_map,
                        prefix=f"{shape_id}_",
                    )
                except Exception:
                    pass
                continue

            # Chart
            if hasattr(shape, "has_chart") and shape.has_chart:
                self._rebuild_chart(shape.chart, block_id_base, block_map)

            if shape.has_table:
                self._rebuild_table(shape.table, block_id_base, group_map)

            if shape.has_text_frame:
                self._rebuild_text_frame(
                    shape.text_frame, block_id_base, group_map
                )
            elif not shape.has_table and not (
                hasattr(shape, "has_chart") and shape.has_chart
            ):
                # Fallback shape (no text_frame). Single block keyed by id.
                fallback = block_map.get(block_id_base)
                if fallback is not None:
                    translated = self._best_text(fallback)
                    if translated and translated != fallback.source_text:
                        try:
                            shape.text = _flatten_for_paragraph(translated)
                        except Exception:
                            logger.debug(
                                "Could not write fallback shape %s",
                                block_id_base,
                            )

    def _rebuild_text_frame(
        self,
        tf,
        group_id: str,
        group_map: dict[str, dict[int, ContentBlock]],
    ) -> None:
        paras_blocks = group_map.get(group_id)
        if not paras_blocks:
            return
        joined_src, joined_tgt = write_paragraphs_by_index(tf, paras_blocks)
        if not joined_tgt:
            return
        # Aggregate-level overflow warning (informational only).
        for block in paras_blocks.values():
            warn_overflow(
                block.source_text,
                self._best_text(block) or "",
                block.id,
            )
        enable_autofit(tf)
        adjust_runs_font_size(tf, joined_src, joined_tgt)

    def _rebuild_chart(
        self, chart, block_id_base: str, block_map: dict[str, ContentBlock]
    ) -> None:
        """Write translated text back into chart elements."""
        title_block = block_map.get(f"{block_id_base}_chart_title")
        if title_block:
            translated = self._best_text(title_block)
            if translated and translated != title_block.source_text:
                try:
                    if chart.has_title and chart.chart_title.has_text_frame:
                        title_tf = chart.chart_title.text_frame
                        # Chart titles are typically a single paragraph; write
                        # the translation into the first paragraph's first run.
                        flat = _flatten_for_paragraph(translated)
                        paras = list(title_tf.paragraphs)
                        if paras and paras[0].runs:
                            paras[0].runs[0].text = flat
                            for r in paras[0].runs[1:]:
                                r.text = ""
                except Exception:
                    logger.debug("Could not rebuild chart title %s", block_id_base)

        cats_block = block_map.get(f"{block_id_base}_chart_cats")
        if cats_block:
            translated = self._best_text(cats_block)
            if translated and translated != cats_block.source_text:
                logger.info(
                    "Chart category labels translated but chart data source "
                    "may need manual update: %s", block_id_base
                )

    def _rebuild_table(
        self,
        table,
        block_id_base: str,
        group_map: dict[str, dict[int, ContentBlock]],
    ) -> None:
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                cell_group = f"{block_id_base}_r{row_idx}c{col_idx}"
                paras_blocks = group_map.get(cell_group)
                if not paras_blocks:
                    continue
                joined_src, joined_tgt = write_paragraphs_by_index(
                    cell.text_frame, paras_blocks
                )
                if not joined_tgt:
                    continue
                for block in paras_blocks.values():
                    warn_overflow(
                        block.source_text,
                        self._best_text(block) or "",
                        block.id,
                    )
                enable_autofit(cell.text_frame)
                adjust_runs_font_size(cell.text_frame, joined_src, joined_tgt)

    def _rebuild_notes(
        self,
        slide,
        slide_idx: int,
        group_map: dict[str, dict[int, ContentBlock]],
    ) -> None:
        notes_group = f"slide{slide_idx}_note"
        paras_blocks = group_map.get(notes_group)
        if not paras_blocks:
            return
        if not slide.has_notes_slide:
            return
        notes_tf = slide.notes_slide.notes_text_frame
        write_paragraphs_by_index(notes_tf, paras_blocks)
