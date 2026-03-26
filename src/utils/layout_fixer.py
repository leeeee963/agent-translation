"""PPT layout fixer: auto-fit text frames and proportional font size adjustment.

Applied during PPTX rebuild to prevent text overflow after translation.
"""

from __future__ import annotations

import logging

from pptx.oxml.ns import qn
from pptx.util import Pt

try:
    import lxml.etree as etree
except ImportError:
    etree = None  # type: ignore

logger = logging.getLogger(__name__)

_MIN_FONT_PT = 8.0
_TRIGGER_RATIO = 1.4  # start reducing when translated text is 40% longer


def enable_autofit(text_frame) -> None:
    """Enable normAutofit on a PPT text frame so PPT auto-shrinks font to fit.

    Preserves the original fontScale value from any existing normAutofit element.
    This prevents translated text from appearing oversized when it is shorter than
    the original: PowerPoint stored fontScale=N% because the original text needed
    shrinking; if we discard that and the new text is shorter, PowerPoint recalculates
    to fontScale=100% and the font jumps back to its raw (large) stored size.
    """
    if etree is None:
        return
    try:
        txBody = text_frame._txBody
        for bodyPr in txBody.iter(qn("a:bodyPr")):
            # Preserve fontScale from existing normAutofit before removing it
            existing_font_scale = None
            for el in bodyPr.findall(qn("a:normAutofit")):
                existing_font_scale = el.get("fontScale")
                bodyPr.remove(el)
            for tag in ("a:noAutofit", "a:spAutoFit"):
                for el in bodyPr.findall(qn(tag)):
                    bodyPr.remove(el)
            new_fit = etree.SubElement(bodyPr, qn("a:normAutofit"))
            if existing_font_scale:
                new_fit.set("fontScale", existing_font_scale)
    except Exception as e:
        logger.debug("enable_autofit failed: %s", e)


def adjust_runs_font_size(
    text_frame,
    source_text: str,
    translated_text: str,
) -> None:
    """Proportionally reduce run font sizes when translated text is significantly longer.

    This is a fallback for when normAutofit alone would shrink text too small.
    """
    if not source_text or not translated_text:
        return
    ratio = len(translated_text) / len(source_text)
    if ratio <= _TRIGGER_RATIO:
        return
    try:
        for para in text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    current_pt = run.font.size.pt
                    new_pt = max(current_pt / ratio, _MIN_FONT_PT)
                    if new_pt < current_pt:
                        run.font.size = Pt(new_pt)
    except Exception as e:
        logger.debug("adjust_runs_font_size failed: %s", e)
