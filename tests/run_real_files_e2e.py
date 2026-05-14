"""End-to-end structural validation against two real files.

Runs the full parse → segment → (simulated translation) → rebuild pipeline on
real PPTX and DOCX files. The "translation" is a deterministic [T]-prefix so
we can verify every original paragraph received the right content at rebuild
time without burning real LLM credits. The format-preservation fix is what's
under test; LLM quality is a separate concern.
"""

from __future__ import annotations

import shutil
import sys
from pathlib import Path
from collections import Counter

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from docx import Document
from pptx import Presentation

from src.parser.docx_parser import DocxParser
from src.parser.pptx_parser import PptxParser
from src.translator.agent import (
    _BLOCK_MARKER,
    _build_marked_input,
    _parse_marked_response,
)
from src.translator.segmenter import Segmenter


def _translate_marker(text: str) -> str:
    return f"[T]{text}"


def _faithful_response(blocks):
    parts = []
    last_group = None
    for b in blocks:
        text = _translate_marker(b.source_text)
        marker = f"{_BLOCK_MARKER.format(block_id=b.id)}\n{text}"
        group = b.metadata.get("text_frame_group") if b.metadata else None
        if group and group != last_group:
            parts.append(f"<!-- group: {group} -->\n{marker}")
        else:
            parts.append(marker)
        last_group = group
    return "\n\n".join(parts) + "\n\n[[END]]"


def _apply_translations(parsed, file_type: str):
    seg_engine = Segmenter()
    blocks = parsed.translatable_blocks
    max_tokens = 5200 if file_type == "pptx" else 3800
    segments = seg_engine.segment(blocks, file_type=file_type, max_tokens=max_tokens)
    print(f"  segments: {len(segments)}")
    for seg in segments:
        # Exercise the prompt builder (this is what would be sent to the LLM).
        _ = _build_marked_input(seg)
        # Apply faithful translation.
        response = _faithful_response(seg)
        matched = _parse_marked_response(response, seg)
        if matched != {b.id for b in seg}:
            missing = {b.id for b in seg} - matched
            print(f"  WARN: segment had {len(missing)} unmatched blocks")


# ---------------------------------------------------------------------------
# PPTX verification
# ---------------------------------------------------------------------------


def verify_pptx(src_path: str, out_path: str) -> dict:
    """Verify rebuilt PPTX preserves every original paragraph."""
    src = Presentation(src_path)
    out = Presentation(out_path)

    failures: list[str] = []
    stats = Counter()

    if len(src.slides) != len(out.slides):
        failures.append(
            f"slide count changed: {len(src.slides)} → {len(out.slides)}"
        )

    def walk_shape_paras(shapes, prefix=""):
        out = []
        for i, shape in enumerate(shapes):
            sid = f"{prefix}shape{i}"
            if shape.shape_type == 6:  # group
                out.extend(walk_shape_paras(shape.shapes, prefix=f"{sid}_"))
                continue
            if shape.has_table:
                for r, row in enumerate(shape.table.rows):
                    for c, cell in enumerate(row.cells):
                        out.append(
                            (
                                f"{sid}_r{r}c{c}",
                                [p.text for p in cell.text_frame.paragraphs],
                            )
                        )
            if shape.has_text_frame:
                out.append((sid, [p.text for p in shape.text_frame.paragraphs]))
        return out

    for slide_idx in range(min(len(src.slides), len(out.slides))):
        src_paras = walk_shape_paras(src.slides[slide_idx].shapes)
        out_paras = walk_shape_paras(out.slides[slide_idx].shapes)
        src_map = dict(src_paras)
        out_map = dict(out_paras)

        for shape_id, src_texts in src_map.items():
            out_texts = out_map.get(shape_id, [])
            if len(src_texts) != len(out_texts):
                failures.append(
                    f"slide {slide_idx} {shape_id}: paragraph count "
                    f"{len(src_texts)} → {len(out_texts)}"
                )
                continue
            for p_idx, (src_t, out_t) in enumerate(zip(src_texts, out_texts)):
                stats["paras_total"] += 1
                if not src_t.strip():
                    if out_t.strip():
                        failures.append(
                            f"slide {slide_idx} {shape_id} p{p_idx}: "
                            f"spacer became {out_t!r}"
                        )
                    else:
                        stats["paras_spacer_preserved"] += 1
                    continue
                # Translatable paragraph: should now contain [T]<source> for the
                # text part. Icon-only runs make src_t exactly "<icon ligature>"
                # — for those we expect no change.
                if src_t == out_t:
                    stats["paras_unchanged"] += 1  # likely icon-only or non-translatable
                elif f"[T]{src_t.strip()}" in out_t:
                    stats["paras_translated_strict"] += 1
                elif "[T]" in out_t:
                    stats["paras_translated_loose"] += 1
                else:
                    failures.append(
                        f"slide {slide_idx} {shape_id} p{p_idx}: src "
                        f"{src_t!r} → out {out_t!r} (no [T] prefix)"
                    )

    return {"failures": failures, "stats": dict(stats)}


# ---------------------------------------------------------------------------
# DOCX verification
# ---------------------------------------------------------------------------


def verify_docx(src_path: str, out_path: str) -> dict:
    src = Document(src_path)
    out = Document(out_path)

    failures: list[str] = []
    stats = Counter()

    if len(list(src.paragraphs)) != len(list(out.paragraphs)):
        failures.append(
            f"body paragraph count changed: "
            f"{len(list(src.paragraphs))} → {len(list(out.paragraphs))}"
        )

    for i, (sp, op) in enumerate(zip(src.paragraphs, out.paragraphs)):
        stats["body_paras_total"] += 1
        if "\n" in op.text or "\r" in op.text:
            failures.append(f"body p{i}: literal newline in rebuilt text")
        if not sp.text.strip():
            if op.text.strip():
                failures.append(f"body p{i}: blank → {op.text!r}")
            else:
                stats["body_blank_preserved"] += 1
            continue
        if sp.text == op.text:
            stats["body_unchanged"] += 1  # non-translatable (URL, numbers, short)
        elif "[T]" in op.text:
            stats["body_translated"] += 1
        else:
            failures.append(
                f"body p{i}: {sp.text!r} → {op.text!r} (no [T])"
            )

    if len(src.tables) != len(out.tables):
        failures.append(f"table count changed: {len(src.tables)} → {len(out.tables)}")

    for t_idx, (st, ot) in enumerate(zip(src.tables, out.tables)):
        for r, (srow, orow) in enumerate(zip(st.rows, ot.rows)):
            for c, (sc, oc) in enumerate(zip(srow.cells, orow.cells)):
                if len(sc.paragraphs) != len(oc.paragraphs):
                    failures.append(
                        f"table {t_idx} r{r}c{c}: paragraph count changed"
                    )
                    continue
                for p_idx, (sp_, op_) in enumerate(zip(sc.paragraphs, oc.paragraphs)):
                    stats["cell_paras_total"] += 1
                    if "\n" in op_.text:
                        failures.append(
                            f"table {t_idx} r{r}c{c} p{p_idx}: literal newline"
                        )
                    if not sp_.text.strip():
                        if op_.text.strip():
                            failures.append(
                                f"table {t_idx} r{r}c{c} p{p_idx}: "
                                f"blank → {op_.text!r}"
                            )
                        continue
                    if sp_.text == op_.text:
                        stats["cell_unchanged"] += 1
                    elif "[T]" in op_.text:
                        stats["cell_translated"] += 1
                    else:
                        failures.append(
                            f"table {t_idx} r{r}c{c} p{p_idx}: "
                            f"{sp_.text!r} → {op_.text!r} (no [T])"
                        )

    return {"failures": failures, "stats": dict(stats)}


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def run(src: str, kind: str):
    print(f"\n=== {kind.upper()}: {Path(src).name} ===")
    work_dir = Path("/tmp") / "agent_translation_e2e"
    work_dir.mkdir(exist_ok=True)
    src_copy = work_dir / Path(src).name
    out_path = work_dir / f"out_{Path(src).stem}.{kind}"
    shutil.copy2(src, src_copy)

    parser = PptxParser() if kind == "pptx" else DocxParser()
    print(f"  parsing {src_copy}")
    parsed = parser.parse(str(src_copy))
    print(f"  blocks parsed: {len(parsed.blocks)} (translatable: {len(parsed.translatable_blocks)})")

    groups = {b.metadata.get("text_frame_group") for b in parsed.translatable_blocks}
    print(f"  text_frame_groups: {len(groups - {None})}")

    _apply_translations(parsed, file_type=kind)

    print(f"  rebuilding → {out_path}")
    parser.rebuild(parsed, str(out_path))

    print(f"  verifying...")
    if kind == "pptx":
        result = verify_pptx(str(src_copy), str(out_path))
    else:
        result = verify_docx(str(src_copy), str(out_path))

    print(f"  stats: {result['stats']}")
    failures = result["failures"]
    if failures:
        print(f"  FAILURES ({len(failures)}):")
        for f in failures[:20]:
            print(f"    • {f}")
        if len(failures) > 20:
            print(f"    … and {len(failures) - 20} more")
        return False
    print(f"  ✓ all assertions passed; output: {out_path}")
    return True


if __name__ == "__main__":
    pptx = "/Users/lee/Downloads/DSCI 5800 Week 2 Lectures.pptx"
    docx = "/Users/lee/Downloads/DSCI 5900 Introduction to AI and Data Science ProgrammingSyllabus .docx"
    p_ok = run(pptx, "pptx")
    d_ok = run(docx, "docx")
    print()
    print("=" * 60)
    print(f"PPTX: {'PASS' if p_ok else 'FAIL'}")
    print(f"DOCX: {'PASS' if d_ok else 'FAIL'}")
    sys.exit(0 if (p_ok and d_ok) else 1)
