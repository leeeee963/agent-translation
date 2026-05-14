"""Adversarial test harness for the per-paragraph PPTX translation pipeline.

Goal: prove that no LLM misbehavior can produce silent paragraph-loss or
paragraph-merge in the rebuilt PPTX. We build synthetic .pptx files that
exercise every shape kind we care about, mock the LLM with progressively
nastier responses, then assert the rebuilt PPTX paragraph-by-paragraph.

Run with:  .venv/bin/python -m pytest tests/test_pptx_per_paragraph.py -v
"""

from __future__ import annotations

import re
import tempfile
from collections.abc import Callable
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from src.models.content import ContentBlock, ParsedFile
from src.parser.pptx_parser import PptxParser
from src.translator.agent import (
    _BLOCK_MARKER,
    _build_marked_input,
    _parse_marked_response,
)
from src.translator.segmenter import Segmenter


# ---------------------------------------------------------------------------
# Synthetic PPTX builders
# ---------------------------------------------------------------------------


def _add_text_box(
    slide,
    paragraphs: list[str | None],
    *,
    left=Inches(1),
    top=Inches(1),
    width=Inches(6),
    height=Inches(4),
    bold_first: bool = False,
    first_size_pt: int | None = None,
):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    # tf starts with one empty paragraph; reuse it for paragraph[0]
    for i, text in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if text is None:
            # Spacer — leave the paragraph with no runs.
            continue
        run = p.add_run()
        run.text = text
        if i == 0 and bold_first:
            run.font.bold = True
        if i == 0 and first_size_pt is not None:
            run.font.size = Pt(first_size_pt)
    return tx


def _add_icon_text_box(slide, *, left=Inches(1), top=Inches(1), width=Inches(6), height=Inches(4)):
    """Build a text box that mixes Material Icons runs with plain text.

    Layout:
        Para 0: "Header" (bold)
        Para 1: "" (spacer)
        Para 2: icon-only (Material Icons, "check_circle")
        Para 3: "Real content"
        Para 4: icon-prefixed text — icon run + text run in same paragraph
    """
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame

    # Para 0: Header, bold
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = "Header"
    r0.font.bold = True

    # Para 1: spacer
    tf.add_paragraph()

    # Para 2: icon-only
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "check_circle"
    r2.font.name = "Material Icons"

    # Para 3: real content
    p3 = tf.add_paragraph()
    r3 = p3.add_run()
    r3.text = "Real content"

    # Para 4: icon + text mixed
    p4 = tf.add_paragraph()
    ricon = p4.add_run()
    ricon.text = "favorite"
    ricon.font.name = "Material Icons"
    rtext = p4.add_run()
    rtext.text = " liked item"

    return tx


def _add_table(
    slide,
    cells: list[list[list[str | None]]],
    *,
    left=Inches(1),
    top=Inches(1),
    width=Inches(8),
    height=Inches(4),
):
    rows = len(cells)
    cols = len(cells[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = tbl_shape.table
    for r, row in enumerate(cells):
        for c, paragraphs in enumerate(row):
            cell = table.cell(r, c)
            tf = cell.text_frame
            for i, text in enumerate(paragraphs):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                if text is None:
                    continue
                run = p.add_run()
                run.text = text
    return tbl_shape


def build_adversarial_pptx() -> str:
    """Build a PPTX exercising every translation hazard.

    Slide 1: header + 3 bullets (the original failure case).
    Slide 2: header + spacer + bullets + spacer + bullet.
    Slide 3: 2x2 table with multi-paragraph cells, including an empty cell.
    Slide 4: grouped shape containing two text frames.
    Slide 5: notes with multiple paragraphs.
    Slide 6: text frame whose every paragraph carries text (no spacers).
    """
    prs = Presentation()
    blank = prs.slide_layouts[6]  # blank layout

    # Slide 1
    s1 = prs.slides.add_slide(blank)
    _add_text_box(
        s1,
        [
            "Characteristics",
            "Strong long-context orientation.",
            "Designed for multi-step agent workflows.",
            "Useful for code, documents, and research tasks.",
        ],
        bold_first=True,
        first_size_pt=18,
    )

    # Slide 2
    s2 = prs.slides.add_slide(blank)
    _add_text_box(
        s2,
        [
            "Pros",
            None,  # spacer
            "Good for reading large inputs.",
            None,
            "Useful for decomposing multi-part tasks.",
        ],
        bold_first=True,
    )

    # Slide 3 — table
    s3 = prs.slides.add_slide(blank)
    _add_table(
        s3,
        [
            [
                ["Title", "Subtitle line"],
                ["Single cell content"],
            ],
            [
                ["Alpha", "Beta", "Gamma"],
                [],  # empty cell
            ],
        ],
    )

    # Slide 4 — grouped shape (build two text frames then group)
    s4 = prs.slides.add_slide(blank)
    tx_a = _add_text_box(
        s4,
        ["Group A", "First note A", "Second note A"],
        left=Inches(0.5), top=Inches(1), width=Inches(4), height=Inches(3),
    )
    tx_b = _add_text_box(
        s4,
        ["Group B", "First note B", "Second note B"],
        left=Inches(5), top=Inches(1), width=Inches(4), height=Inches(3),
    )
    # python-pptx exposes group creation in 0.6.21+
    try:
        from pptx.util import Emu  # noqa: F401

        s4.shapes.add_group_shape([tx_a, tx_b])
    except Exception:
        # If grouping API is unavailable, leave them as siblings — we still
        # cover the non-grouped multi-textbox case.
        pass

    # Slide 5 — notes
    s5 = prs.slides.add_slide(blank)
    _add_text_box(s5, ["Slide five body"])
    notes_tf = s5.notes_slide.notes_text_frame
    # First paragraph is already empty; populate it then add more.
    notes_tf.paragraphs[0].add_run().text = "First note paragraph."
    p2 = notes_tf.add_paragraph()
    p2.add_run().text = "Second note paragraph."
    p3 = notes_tf.add_paragraph()
    p3.add_run().text = "Third note paragraph with a longer sentence to vary."

    # Slide 6 — every paragraph populated, no spacers
    s6 = prs.slides.add_slide(blank)
    _add_text_box(
        s6,
        ["Line one", "Line two", "Line three", "Line four"],
    )

    # Slide 7 — icon fonts mixed with text
    s7 = prs.slides.add_slide(blank)
    _add_icon_text_box(s7)

    # Slide 8 — stress: 30-bullet list
    s8 = prs.slides.add_slide(blank)
    _add_text_box(
        s8,
        [f"Bullet number {i:02d}" for i in range(30)],
    )

    out = Path(tempfile.gettempdir()) / "adversarial_test.pptx"
    prs.save(str(out))
    return str(out)


# ---------------------------------------------------------------------------
# Adversarial response generators
# ---------------------------------------------------------------------------


def _faithful_response(blocks: list[ContentBlock], translate: Callable[[str], str]) -> str:
    """Echo the marked input back with each source replaced by translate(source)."""
    parts = []
    last_group: str | None = None
    for b in blocks:
        text = translate(b.source_text)
        marker = f"{_BLOCK_MARKER.format(block_id=b.id)}\n{text}"
        group = b.metadata.get("text_frame_group")
        if group and group != last_group:
            parts.append(f"<!-- group: {group} -->\n{marker}")
        else:
            parts.append(marker)
        last_group = group
    return "\n\n".join(parts) + "\n\n[[END]]"


def adversary_drop_last(
    blocks: list[ContentBlock], translate: Callable[[str], str]
) -> str:
    """Drop the last block in every group."""
    keep = []
    by_group: dict[str | None, list[ContentBlock]] = {}
    for b in blocks:
        by_group.setdefault(b.metadata.get("text_frame_group"), []).append(b)
    drop_ids = {grp[-1].id for grp in by_group.values() if len(grp) > 1}
    for b in blocks:
        if b.id not in drop_ids:
            keep.append(b)
    return _faithful_response(keep, translate)


def adversary_merge_two(
    blocks: list[ContentBlock], translate: Callable[[str], str]
) -> str:
    """Merge the first two blocks of every group into a single marker."""
    parts = []
    skip_ids: set[str] = set()
    by_group: dict[str | None, list[ContentBlock]] = {}
    for b in blocks:
        by_group.setdefault(b.metadata.get("text_frame_group"), []).append(b)
    merge_pairs: dict[str, str] = {}
    for grp in by_group.values():
        if len(grp) >= 2:
            merge_pairs[grp[0].id] = grp[1].source_text
            skip_ids.add(grp[1].id)

    last_group: str | None = None
    for b in blocks:
        if b.id in skip_ids:
            continue
        text = translate(b.source_text)
        if b.id in merge_pairs:
            text = text + " " + translate(merge_pairs[b.id])  # merged inline
        marker = f"{_BLOCK_MARKER.format(block_id=b.id)}\n{text}"
        group = b.metadata.get("text_frame_group")
        if group and group != last_group:
            parts.append(f"<!-- group: {group} -->\n{marker}")
        else:
            parts.append(marker)
        last_group = group
    return "\n\n".join(parts) + "\n\n[[END]]"


def adversary_swap_order(
    blocks: list[ContentBlock], translate: Callable[[str], str]
) -> str:
    """Return blocks in reversed order. Marker-based parsing should still work."""
    return _faithful_response(list(reversed(blocks)), translate)


def adversary_extra_noise(
    blocks: list[ContentBlock], translate: Callable[[str], str]
) -> str:
    """Inject prose preamble, blank lines, stray <!-- ... --> comments, and a
    chatty epilogue AFTER [[END]] that the parser must drop on the floor."""
    parts = ["Here are the translations:\n"]
    for i, b in enumerate(blocks):
        text = translate(b.source_text)
        comment_extra = "<!-- nb: machine output -->" if i % 3 == 0 else ""
        marker = f"{_BLOCK_MARKER.format(block_id=b.id)}\n{text}\n{comment_extra}"
        parts.append(marker)
        parts.append("")  # extra blank
    body = "\n\n".join(parts)
    return body + "\n\n[[END]]\n\nLet me know if you need adjustments. — End."


def adversary_empty_one(
    blocks: list[ContentBlock], translate: Callable[[str], str]
) -> str:
    """Return empty translation for the first block of every group."""
    by_group: dict[str | None, list[ContentBlock]] = {}
    for b in blocks:
        by_group.setdefault(b.metadata.get("text_frame_group"), []).append(b)
    empty_ids = {grp[0].id for grp in by_group.values()}
    parts = []
    last_group: str | None = None
    for b in blocks:
        text = "" if b.id in empty_ids else translate(b.source_text)
        marker = f"{_BLOCK_MARKER.format(block_id=b.id)}\n{text}"
        group = b.metadata.get("text_frame_group")
        if group and group != last_group:
            parts.append(f"<!-- group: {group} -->\n{marker}")
        else:
            parts.append(marker)
        last_group = group
    return "\n\n".join(parts) + "\n\n[[END]]"


def adversary_protocol_injection(
    blocks: list[ContentBlock], translate: Callable[[str], str]
) -> str:
    """Translation text contains literal [[END]] and [[BLOCK:fake_id]].

    A robust parser must require the markers to be on their own line; an
    inline occurrence inside a translation must NOT prematurely truncate or
    confuse block matching.
    """
    parts = []
    last_group: str | None = None
    for i, b in enumerate(blocks):
        text = translate(b.source_text)
        # Sprinkle hostile inline tokens
        if i == 1:
            text = f"{text} (note: [[END]] of section)"
        elif i == 2:
            text = f"{text} — see [[BLOCK:hallucinated]] reference"
        marker = f"{_BLOCK_MARKER.format(block_id=b.id)}\n{text}"
        group = b.metadata.get("text_frame_group")
        if group and group != last_group:
            parts.append(f"<!-- group: {group} -->\n{marker}")
        else:
            parts.append(marker)
        last_group = group
    return "\n\n".join(parts) + "\n\n[[END]]"


def adversary_drop_everything(
    blocks: list[ContentBlock], translate: Callable[[str], str]
) -> str:
    """Returns no block markers at all — every block must come back via the
    single-block fallback."""
    return "[[END]]"


def adversary_hallucinated_ids(
    blocks: list[ContentBlock], translate: Callable[[str], str]
) -> str:
    """Mix real translations with markers for ids that don't exist.

    The fake markers must be silently ignored; the real ones must still match.
    """
    parts = []
    last_group: str | None = None
    for i, b in enumerate(blocks):
        text = translate(b.source_text)
        marker = f"{_BLOCK_MARKER.format(block_id=b.id)}\n{text}"
        group = b.metadata.get("text_frame_group")
        if group and group != last_group:
            parts.append(f"<!-- group: {group} -->\n{marker}")
        else:
            parts.append(marker)
        last_group = group
        # Sprinkle hallucinated markers
        if i % 4 == 0:
            parts.append(
                f"[[BLOCK:totally_fake_id_{i}]]\nthis text should be ignored"
            )
    return "\n\n".join(parts) + "\n\n[[END]]"


def adversary_whitespace_translation(
    blocks: list[ContentBlock], translate: Callable[[str], str]
) -> str:
    """Returns whitespace-only translations for first half of every group.

    These should be treated as missing (per `_parse_marked_response`'s
    `text` truthiness check), triggering retry/fallback.
    """
    by_group: dict[str | None, list[ContentBlock]] = {}
    for b in blocks:
        by_group.setdefault(b.metadata.get("text_frame_group"), []).append(b)
    blank_ids: set[str] = set()
    for grp in by_group.values():
        for b in grp[: max(1, len(grp) // 2)]:
            blank_ids.add(b.id)
    parts = []
    last_group: str | None = None
    for b in blocks:
        text = "   \n  \n   " if b.id in blank_ids else translate(b.source_text)
        marker = f"{_BLOCK_MARKER.format(block_id=b.id)}\n{text}"
        group = b.metadata.get("text_frame_group")
        if group and group != last_group:
            parts.append(f"<!-- group: {group} -->\n{marker}")
        else:
            parts.append(marker)
        last_group = group
    return "\n\n".join(parts) + "\n\n[[END]]"


def adversary_persistent_drop(
    blocks: list[ContentBlock], translate: Callable[[str], str]
) -> str:
    """Always drops the last block of every group, even on retry.

    Forces the single-block fallback path. The pipeline must still produce
    a complete translation for every paragraph.
    """
    by_group: dict[str | None, list[ContentBlock]] = {}
    for b in blocks:
        by_group.setdefault(b.metadata.get("text_frame_group"), []).append(b)
    drop_ids = {grp[-1].id for grp in by_group.values() if len(grp) > 1}
    keep = [b for b in blocks if b.id not in drop_ids]
    return _faithful_response(keep, translate)


def adversary_multiline_for_single(
    blocks: list[ContentBlock], translate: Callable[[str], str]
) -> str:
    """Return multi-line translation for a per-paragraph block. Rebuild must
    flatten the newlines so the source paragraph stays as a single paragraph
    in the output (no ghost paragraphs grown from injected newlines).

    Content-preserving: the adversary only injects whitespace; the flatten
    pass should collapse it back so the test can still check content equality.
    """
    parts = []
    last_group: str | None = None
    for b in blocks:
        text = translate(b.source_text)
        # Inject newlines around/inside without adding tokens.
        if " " in text:
            text = text.replace(" ", "\n", 1)  # split first space → newline
        text = f"\n{text}\n"  # leading + trailing blank lines
        marker = f"{_BLOCK_MARKER.format(block_id=b.id)}\n{text}"
        group = b.metadata.get("text_frame_group")
        if group and group != last_group:
            parts.append(f"<!-- group: {group} -->\n{marker}")
        else:
            parts.append(marker)
        last_group = group
    return "\n\n".join(parts) + "\n\n[[END]]"


# ---------------------------------------------------------------------------
# Test harness
# ---------------------------------------------------------------------------


def _translate_marker(text: str) -> str:
    """Deterministic 'translation' that prefixes [T] so we can assert which
    paragraph received which source and detect cross-talk."""
    return f"[T]{text}"


def _apply_translations(
    parsed: ParsedFile,
    adversary: Callable[[list[ContentBlock], Callable[[str], str]], str],
    *,
    max_retries: int = 2,
    persistent: bool = False,
):
    """Run the per-segment marker pipeline the agent uses, but in-process.

    Mirrors TranslatorAgent._translate_segment's retry behaviour:
    - One LLM call (the adversary).
    - If markers missing, one retry. With ``persistent=True`` the retry uses
      the SAME adversary (simulating an LLM that keeps making the same
      mistake); otherwise the retry is faithful.
    - If still missing, fall back to single-block calls (always faithful —
      single-block calls in production retry independently).
    """
    seg_engine = Segmenter()
    blocks = parsed.translatable_blocks
    segments = seg_engine.segment(blocks, file_type="pptx", max_tokens=5200)

    for seg in segments:
        _ = _build_marked_input(seg)  # exercise the group-marker builder
        response = adversary(seg, _translate_marker)
        matched = _parse_marked_response(response, seg)
        attempts = 1
        while attempts <= max_retries and len(matched) < len(seg):
            for b in seg:
                if b.id not in matched:
                    b.translated_text = ""
            response = (
                adversary(seg, _translate_marker)
                if persistent
                else _faithful_response(seg, _translate_marker)
            )
            matched = _parse_marked_response(response, seg)
            attempts += 1
        if len(matched) < len(seg):
            # Final fallback: single-block translations for EVERY block,
            # mirroring TranslatorAgent._translate_segment. Even matched
            # blocks may carry polluted content from a merge — they must be
            # re-translated, not preserved.
            for b in seg:
                b.translated_text = ""
                single_resp = _faithful_response([b], _translate_marker)
                _parse_marked_response(single_resp, [b])


# ---------------------------------------------------------------------------
# Per-paragraph rebuild assertion helpers
# ---------------------------------------------------------------------------


def _collect_text_frame_paragraphs(prs) -> list[tuple[str, list[str]]]:
    """Walk the PPTX and return [(label, [para_text, ...])] for every text frame
    we care about (text boxes, table cells, notes), preserving order."""
    out: list[tuple[str, list[str]]] = []

    def walk_shapes(shapes, slide_idx: int, prefix: str = ""):
        for shape_idx, shape in enumerate(shapes):
            shape_id = f"{prefix}shape{shape_idx}" if prefix else f"shape{shape_idx}"
            label_base = f"slide{slide_idx}_{shape_id}"
            if shape.shape_type is not None and shape.shape_type == 6:
                walk_shapes(shape.shapes, slide_idx, prefix=f"{shape_id}_")
                continue
            if shape.has_table:
                for r, row in enumerate(shape.table.rows):
                    for c, cell in enumerate(row.cells):
                        out.append(
                            (
                                f"{label_base}_r{r}c{c}",
                                [p.text for p in cell.text_frame.paragraphs],
                            )
                        )
            if shape.has_text_frame:
                out.append(
                    (label_base, [p.text for p in shape.text_frame.paragraphs])
                )

    for slide_idx, slide in enumerate(prs.slides):
        walk_shapes(slide.shapes, slide_idx)
        if slide.has_notes_slide:
            ntf = slide.notes_slide.notes_text_frame
            out.append((f"slide{slide_idx}_note", [p.text for p in ntf.paragraphs]))

    return out


def _expected_after_translation(parsed: ParsedFile) -> dict[str, dict[int, str]]:
    """Expected per-paragraph text after a faithful '[T]'-prefix translation.

    Mirrors what ``write_para_with_fmt`` does at rebuild time:
      • icon-font runs in the original keep their text in place
      • the first non-icon run becomes ``lead_ws + [T]<source> + trail_ws``
      • any other non-icon runs are cleared

    So the expected ``para.text`` for an icon-prefixed paragraph like
    ``<icon>favorite</icon> liked item`` is ``"favorite [T]liked item"``.
    """
    # {group: {(para_idx, line_idx): expected_full_para_text_contribution}}
    # We collapse to {group: {para_idx: concatenated text}} because para.text
    # is the visible string at verification time.
    by_group_para: dict[str, dict[int, list[tuple[int, str]]]] = {}
    for block in parsed.blocks:
        group = block.metadata.get("text_frame_group")
        para_idx = block.metadata.get("para_index")
        if group is None or para_idx is None:
            continue
        line_idx = int(block.metadata.get("line_index") or 0)
        line_fmt = block.metadata.get("line_format") or {}
        runs_kind = line_fmt.get("runs_kind") or []
        trailing_break = line_fmt.get("trailing_break")
        lead_ws = block.metadata.get("lead_ws", "")
        trail_ws = block.metadata.get("trail_ws", "")
        text_part = lead_ws + f"[T]{block.source_text}" + trail_ws
        if not runs_kind:
            line_str = text_part
        else:
            parts: list[str] = []
            written = False
            for run in runs_kind:
                if run.get("is_icon"):
                    parts.append(run.get("text", ""))
                elif not written:
                    parts.append(text_part)
                    written = True
            if not written:
                parts.append(text_part)
            line_str = "".join(parts)
        if trailing_break == "newline":
            line_str = line_str + "\n"
        by_group_para.setdefault(group, {}).setdefault(int(para_idx), []).append(
            (line_idx, line_str)
        )

    expected: dict[str, dict[int, str]] = {}
    for group, paras in by_group_para.items():
        for para_idx, lines in paras.items():
            lines.sort(key=lambda kv: kv[0])
            expected.setdefault(group, {})[para_idx] = "".join(s for _, s in lines)
    return expected


def _assert_rebuild_matches(
    output_path: str,
    parsed_source: ParsedFile,
    source_pptx: str,
    *,
    strict_content: bool = True,
):
    """Assert the rebuilt PPTX matches expectations.

    Two property layers are checked:

    1. STRUCTURAL (always strict):
       • Every block-paragraph receives content at its (group, para_index).
       • Every non-block paragraph (spacer, icon-only) is byte-identical to
         the source — the rebuild must not touch them.
       • No paragraph dropped, no paragraph created.

    2. CONTENT (``strict_content=True`` only): each block-paragraph's rebuilt
       text equals ``[T]<source>`` exactly. Set ``False`` for adversaries
       that mutate translation content.
    """
    expected = _expected_after_translation(parsed_source)
    original_paras = dict(_collect_text_frame_paragraphs(Presentation(source_pptx)))
    actual_paras = dict(_collect_text_frame_paragraphs(Presentation(output_path)))

    failures: list[str] = []
    for label, original in original_paras.items():
        actual = actual_paras.get(label, [])
        if len(actual) != len(original):
            failures.append(
                f"{label}: paragraph count changed "
                f"(orig {len(original)}, rebuilt {len(actual)})"
            )
        expected_paras = expected.get(label, {})
        for para_idx, orig_text in enumerate(original):
            if para_idx >= len(actual):
                continue
            got = actual[para_idx]
            if para_idx in expected_paras:
                # Block paragraph — must contain the translated content.
                expected_text = expected_paras[para_idx]
                if strict_content:
                    if got.strip() != expected_text:
                        failures.append(
                            f"{label} para {para_idx}: expected {expected_text!r}, got {got!r}"
                        )
                else:
                    if expected_text not in got:
                        failures.append(
                            f"{label} para {para_idx}: missing source content "
                            f"{expected_text!r} (got {got!r})"
                        )
            else:
                # Non-block paragraph (spacer or icon-only) — must be identical.
                if got != orig_text:
                    failures.append(
                        f"{label} para {para_idx} was non-translatable but changed: "
                        f"{orig_text!r} → {got!r}"
                    )

    if failures:
        raise AssertionError(
            "Rebuild mismatch:\n  " + "\n  ".join(failures)
        )


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------


@pytest.fixture(scope="module")
def adversarial_pptx_path() -> str:
    return build_adversarial_pptx()


# content_mutating=True means the adversary deliberately adds extra characters
# to the translation; only the structural integrity is asserted in that case.
@pytest.mark.parametrize(
    "adversary_name,adversary_fn,persistent,content_mutating",
    [
        ("faithful", _faithful_response, False, False),
        ("drop_last", adversary_drop_last, False, False),
        ("merge_two", adversary_merge_two, False, False),
        ("swap_order", adversary_swap_order, False, False),
        ("extra_noise", adversary_extra_noise, False, False),
        ("empty_one", adversary_empty_one, False, False),
        ("multiline", adversary_multiline_for_single, False, False),
        ("protocol_injection", adversary_protocol_injection, False, True),
        # Persistent: same adversary fires on every retry; only the
        # single-block fallback prevents file corruption.
        ("persistent_drop_last", adversary_persistent_drop, True, False),
        ("persistent_merge_two", adversary_merge_two, True, False),
        ("persistent_empty_one", adversary_empty_one, True, False),
        ("persistent_protocol_injection", adversary_protocol_injection, True, True),
        ("drop_everything", adversary_drop_everything, False, False),
        ("hallucinated_ids", adversary_hallucinated_ids, False, False),
        ("whitespace_translation", adversary_whitespace_translation, False, False),
        ("persistent_drop_everything", adversary_drop_everything, True, False),
        ("persistent_whitespace", adversary_whitespace_translation, True, False),
    ],
)
def test_adversarial_rebuild(
    adversarial_pptx_path,
    adversary_name,
    adversary_fn,
    persistent,
    content_mutating,
    tmp_path,
):
    parser = PptxParser()
    parsed = parser.parse(adversarial_pptx_path)

    _apply_translations(parsed, adversary_fn, persistent=persistent)

    out_path = str(tmp_path / f"out_{adversary_name}.pptx")
    parser.rebuild(parsed, out_path)

    _assert_rebuild_matches(
        out_path,
        parsed_source=parsed,
        source_pptx=adversarial_pptx_path,
        strict_content=not content_mutating,
    )


def test_per_paragraph_block_ids(adversarial_pptx_path):
    """Sanity: every translatable text-frame block has the per-paragraph schema."""
    parsed = PptxParser().parse(adversarial_pptx_path)
    text_frame_blocks = [
        b for b in parsed.translatable_blocks
        if b.metadata.get("group_kind") in {"text_frame", "table_cell", "notes"}
    ]
    assert text_frame_blocks, "expected per-paragraph text frame blocks"
    for b in text_frame_blocks:
        assert b.metadata.get("text_frame_group"), b.id
        assert b.metadata.get("para_index") is not None, b.id
        assert b.metadata.get("line_index") is not None, b.id
        # Block id format: ..._p{K}_l{L} (paragraph + soft-line within).
        assert re.search(r"_p\d+_l\d+$", b.id), b.id
        # Source text is a single soft-line — must not contain newlines.
        assert "\n" not in b.source_text, b.id


def test_group_markers_in_prompt(adversarial_pptx_path):
    """The prompt input must group blocks of the same text_frame_group under one
    `<!-- group: X -->` comment so the LLM can see them as a unit."""
    parsed = PptxParser().parse(adversarial_pptx_path)
    blocks = parsed.translatable_blocks
    msg = _build_marked_input(blocks)

    # Every text_frame_group should appear exactly once as a group comment.
    groups = {b.metadata.get("text_frame_group") for b in blocks if b.metadata.get("text_frame_group")}
    for g in groups:
        assert msg.count(f"<!-- group: {g} -->") == 1, (
            f"group {g!r} not introduced exactly once in prompt input"
        )


def test_icon_fonts_are_preserved(adversarial_pptx_path, tmp_path):
    """Icon-font runs must (a) never be sent to the LLM, (b) never be touched
    on rebuild. Their ligature glyphs render correctly only with the original
    font; CJK fallback would show squares or wrong-style characters."""
    parser = PptxParser()
    parsed = parser.parse(adversarial_pptx_path)

    # Sanity: no block was created for the icon-only paragraph.
    icon_only_blocks = [
        b for b in parsed.translatable_blocks
        if b.source_text.strip() == "check_circle"
    ]
    assert not icon_only_blocks, (
        "icon-only paragraph leaked into translation: "
        f"{[b.id for b in icon_only_blocks]}"
    )

    # Sanity: mixed icon+text paragraph carries only the text, not the icon.
    mixed_blocks = [
        b for b in parsed.translatable_blocks
        if "liked item" in b.source_text
    ]
    assert mixed_blocks, "expected the mixed icon+text paragraph as a block"
    for b in mixed_blocks:
        assert "favorite" not in b.source_text, (
            f"icon ligature 'favorite' leaked into block {b.id}: {b.source_text!r}"
        )

    # Translate everything faithfully and rebuild.
    _apply_translations(parsed, _faithful_response)
    out = str(tmp_path / "out_icons.pptx")
    parser.rebuild(parsed, out)

    # Walk the rebuilt file: the icon-only paragraph must still say "check_circle"
    # in a Material Icons run; the mixed paragraph must still have its
    # Material Icons run intact.
    prs = Presentation(out)
    icon_only_seen = False
    mixed_seen = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    fname = run.font.name or ""
                    if fname == "Material Icons" and run.text == "check_circle":
                        icon_only_seen = True
                    if fname == "Material Icons" and run.text == "favorite":
                        mixed_seen = True
    assert icon_only_seen, "icon-only paragraph lost its Material Icons run"
    assert mixed_seen, "mixed paragraph lost its Material Icons run"


def test_format_preserved_after_translation(adversarial_pptx_path, tmp_path):
    """Bold/size on a header paragraph must survive translation."""
    parser = PptxParser()
    parsed = parser.parse(adversarial_pptx_path)
    _apply_translations(parsed, _faithful_response)
    out = str(tmp_path / "out_fmt.pptx")
    parser.rebuild(parsed, out)

    prs = Presentation(out)
    s1 = prs.slides[0]  # "Characteristics" header (bold, 18pt) + 3 bullets
    tf = next(s for s in s1.shapes if s.has_text_frame).text_frame
    header_para = tf.paragraphs[0]
    header_run = header_para.runs[0]
    assert header_run.text.startswith("[T]Characteristics"), header_run.text
    assert header_run.font.bold is True, "header paragraph lost bold"
    assert header_run.font.size == Pt(18), (
        f"header paragraph lost font size: {header_run.font.size}"
    )


def test_multi_segment_with_mixed_adversaries(adversarial_pptx_path, tmp_path):
    """Force the segmenter to produce many small segments and assign a
    different (often nasty) adversary to each. The pipeline must still
    deliver every paragraph correctly."""
    parser = PptxParser()
    parsed = parser.parse(adversarial_pptx_path)

    blocks = parsed.translatable_blocks
    # Force many tiny segments via group-atomic splitting.
    segments = Segmenter()._segment_by_groups(blocks, max_tokens=15)
    assert len(segments) >= 5, f"expected many segments, got {len(segments)}"

    adversaries = [
        _faithful_response,
        adversary_drop_last,
        adversary_merge_two,
        adversary_empty_one,
        adversary_whitespace_translation,
        adversary_hallucinated_ids,
        adversary_protocol_injection,
        adversary_drop_everything,
    ]

    for i, seg in enumerate(segments):
        adv = adversaries[i % len(adversaries)]
        response = adv(seg, _translate_marker)
        matched = _parse_marked_response(response, seg)
        retries = 0
        while retries < 2 and len(matched) < len(seg):
            for b in seg:
                if b.id not in matched:
                    b.translated_text = ""
            response = _faithful_response(seg, _translate_marker)
            matched = _parse_marked_response(response, seg)
            retries += 1
        if len(matched) < len(seg):
            for b in seg:
                b.translated_text = ""
                _parse_marked_response(_faithful_response([b], _translate_marker), [b])

    out = str(tmp_path / "out_multi_segment.pptx")
    parser.rebuild(parsed, out)
    _assert_rebuild_matches(
        out, parsed_source=parsed, source_pptx=adversarial_pptx_path,
        strict_content=False,  # protocol_injection mutates content
    )


def test_original_screenshot_scenario(tmp_path):
    """Reproduce the exact failure mode from the user's report: a text box
    with a bold header followed by 3 bullets. Old code dropped the last bullet
    or merged the header with bullet 1. Per-paragraph blocks must fix it."""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_text_box(
        s,
        [
            "Characteristics",
            "Strong long-context orientation.",
            "Designed for multi-step agent workflows.",
            "Useful for code, documents, and research tasks.",
        ],
        bold_first=True,
        first_size_pt=18,
    )
    src = str(tmp_path / "screenshot.pptx")
    prs.save(src)

    parser = PptxParser()
    parsed = parser.parse(src)

    # Simulate the bad LLM behavior from the screenshot: dropped last bullet
    # (image 2) and merged-header-with-first-bullet (image 1). Pipeline must
    # recover via retries / single-block fallback.
    _apply_translations(parsed, adversary_drop_last, persistent=True)

    out = str(tmp_path / "out_screenshot.pptx")
    parser.rebuild(parsed, out)
    _assert_rebuild_matches(out, parsed_source=parsed, source_pptx=src)


def test_end_sentinel_collision(tmp_path):
    """If a translation legitimately contains '[[END]]' on its own line — for
    example translating documentation about the system itself — the parser
    must NOT prematurely truncate and lose subsequent block markers."""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_text_box(s, ["First paragraph", "Second paragraph", "Third paragraph"])
    src = str(tmp_path / "end_collision.pptx")
    prs.save(src)

    parser = PptxParser()
    parsed = parser.parse(src)
    blocks = parsed.translatable_blocks
    assert len(blocks) == 3

    # Craft an adversary response where block 1's translation contains a line
    # that is exactly "[[END]]" surrounded by [T]-prefixed legitimate content.
    response = (
        f"[[BLOCK:{blocks[0].id}]]\n"
        f"[T]{blocks[0].source_text}\n\n"
        f"[[BLOCK:{blocks[1].id}]]\n"
        f"[T]{blocks[1].source_text}\n"
        f"[[END]]\n"  # legitimate content that LOOKS like the sentinel
        f"continuation line\n\n"
        f"[[BLOCK:{blocks[2].id}]]\n"
        f"[T]{blocks[2].source_text}\n\n"
        f"[[END]]"  # real terminator
    )
    matched = _parse_marked_response(response, blocks)
    # All 3 blocks must match — the parser must use the LAST [[END]] as the
    # terminator, not the first, so the inline one inside block 1 doesn't
    # truncate the response.
    assert matched == {b.id for b in blocks}, (
        f"end-sentinel collision dropped blocks: matched={matched}"
    )

    out = str(tmp_path / "out_end_collision.pptx")
    parser.rebuild(parsed, out)
    _assert_rebuild_matches(
        out, parsed_source=parsed, source_pptx=src,
        strict_content=False,  # block 1 contains extra legitimate content
    )


def test_reviewed_text_takes_precedence(adversarial_pptx_path, tmp_path):
    """When both translated_text and reviewed_text are set, rebuild must use
    the reviewed version (post-naturalness review wins over draft)."""
    parser = PptxParser()
    parsed = parser.parse(adversarial_pptx_path)
    _apply_translations(parsed, _faithful_response)
    # Now simulate a review pass that rewrote everything with a different prefix.
    for b in parsed.translatable_blocks:
        if b.translated_text:
            b.reviewed_text = b.translated_text.replace("[T]", "[R]")
    out = str(tmp_path / "out_reviewed.pptx")
    parser.rebuild(parsed, out)
    prs = Presentation(out)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                t = p.text
                # Where there's translated content, it must be the reviewed one.
                if "[T]" in t and "[R]" not in t:
                    raise AssertionError(
                        f"reviewed_text was not honored: {t!r}"
                    )


def test_soft_line_break_split(tmp_path):
    """Reproduce the *exact* failure pattern from the user's DSCI 5800 deck:
    a single <a:p> contains a bold header + literal '\\n' + a non-bold first
    bullet (a PowerPoint soft line break). Two subsequent <a:p>s carry the
    other bullets.

    Old behavior: header + first bullet land in one block; LLM reshuffles and
    silently drops the last bullet.

    New behavior: each soft line becomes its own block; faithful translation
    preserves all 4 visual lines (header + 3 bullets); the bold/non-bold
    format is preserved per-line; the soft line break is preserved in XML.
    """
    from pptx import Presentation as _P
    from pptx.util import Inches, Pt
    from lxml import etree

    prs = _P()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    tx = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(4))
    tf = tx.text_frame
    # Build paragraph 0: bold header + soft \n + non-bold first bullet
    p0 = tf.paragraphs[0]
    r_hdr = p0.add_run()
    r_hdr.text = "Characteristics\n"  # trailing \n = soft line break
    r_hdr.font.bold = True
    r_bul1 = p0.add_run()
    r_bul1.text = "• Strong long-context orientation."
    # Paragraphs 1 and 2: the other two bullets
    p1 = tf.add_paragraph()
    p1.add_run().text = "• Designed for multi-step agent workflows."
    p2 = tf.add_paragraph()
    p2.add_run().text = "• Useful for code, documents, and research tasks."

    src = str(tmp_path / "softbreak.pptx")
    prs.save(src)

    # --- Parse ---
    parsed = PptxParser().parse(src)
    blocks = sorted(parsed.blocks, key=lambda b: b.id)

    # Expect 4 blocks total: header (p0_l0), first bullet (p0_l1), second
    # bullet (p1_l0), third bullet (p2_l0). NOT 3 blocks.
    assert len(blocks) == 4, f"expected 4 blocks, got {len(blocks)}: {[b.id for b in blocks]}"

    ids = [b.id for b in blocks]
    assert ids[0].endswith("_p0_l0")
    assert ids[1].endswith("_p0_l1")
    assert ids[2].endswith("_p1_l0")
    assert ids[3].endswith("_p2_l0")

    # Source text must not contain newlines (each block is one soft-line).
    for b in blocks:
        assert "\n" not in b.source_text, b.id

    # Per-line format: header is bold, bullets are not.
    fmt0 = blocks[0].metadata["line_format"]
    fmt1 = blocks[1].metadata["line_format"]
    assert fmt0["bold"] is True, "header line lost its bold"
    assert fmt1["bold"] is not True, "bullet line wrongly inherited bold"

    # The soft break must be recorded so rebuild can restore it.
    assert fmt0["trailing_break"] == "newline"
    assert fmt1["trailing_break"] is None

    # --- Faithful translate + rebuild ---
    _apply_translations(parsed, _faithful_response)
    out_path = str(tmp_path / "out_softbreak.pptx")
    PptxParser().rebuild(parsed, out_path)

    # Verify the rebuilt file:
    #  • 3 <a:p> paragraphs preserved (no extra, no missing)
    #  • Paragraph 0 still has the soft line break ('\n' at end of first run)
    #  • All 4 visual lines have their [T]<source> translation
    out = _P(out_path)
    out_tf = list(out.slides[0].shapes)[0].text_frame
    out_paras = list(out_tf.paragraphs)
    assert len(out_paras) == 3, f"paragraph count changed: {len(out_paras)}"

    # Paragraph 0: must contain BOTH the header translation AND the first
    # bullet translation, with a \n between them (the soft line break).
    p0_text = out_paras[0].text
    assert "[T]Characteristics" in p0_text, p0_text
    assert "[T]• Strong long-context orientation." in p0_text, p0_text
    assert "\n" in p0_text, f"soft line break lost in p0: {p0_text!r}"

    # Subsequent bullets in their own paragraphs.
    assert "[T]• Designed for multi-step agent workflows." in out_paras[1].text
    assert "[T]• Useful for code, documents, and research tasks." in out_paras[2].text

    # Format: header run still bold, bullet runs not bold.
    p0_runs = list(out_paras[0].runs)
    hdr_run = next(r for r in p0_runs if "Characteristics" in r.text)
    bul_run = next(r for r in p0_runs if "Strong" in r.text)
    assert hdr_run.font.bold is True, "rebuilt header run lost bold"
    # Bullet run inherits None (theme default) — must not be True.
    assert bul_run.font.bold is not True, "rebuilt bullet run wrongly bold"


def test_soft_line_break_dropped_bullet_resists_adversary(tmp_path):
    """Same source as above, but with an adversary that always drops the last
    block of every group. The pipeline must still recover all 4 lines."""
    from pptx import Presentation as _P
    from pptx.util import Inches

    prs = _P()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tx = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(4))
    tf = tx.text_frame
    p0 = tf.paragraphs[0]
    r_hdr = p0.add_run()
    r_hdr.text = "Characteristics\n"
    r_hdr.font.bold = True
    p0.add_run().text = "• Strong long-context orientation."
    tf.add_paragraph().add_run().text = "• Designed for multi-step agent workflows."
    tf.add_paragraph().add_run().text = "• Useful for code, documents, and research tasks."
    src = str(tmp_path / "softbreak_adv.pptx")
    prs.save(src)

    parser = PptxParser()
    parsed = parser.parse(src)
    _apply_translations(parsed, adversary_drop_last, persistent=True)

    out = str(tmp_path / "out_softbreak_adv.pptx")
    parser.rebuild(parsed, out)
    _assert_rebuild_matches(out, parsed_source=parsed, source_pptx=src)


def test_segmenter_keeps_groups_together(adversarial_pptx_path):
    parsed = PptxParser().parse(adversarial_pptx_path)
    blocks = parsed.translatable_blocks
    # Force a tiny budget so segmenter MUST split.
    segments = Segmenter().segment(blocks, file_type="pptx", max_tokens=20)
    seen_groups: set[str] = set()
    for seg in segments:
        groups_in_seg = {
            b.metadata.get("text_frame_group")
            for b in seg
            if b.metadata.get("text_frame_group")
        }
        for g in groups_in_seg:
            assert g not in seen_groups, (
                f"group {g!r} split across multiple segments — atomicity broken"
            )
            seen_groups.add(g)
